"""Report export rendering — branded, detailed Markdown / PDF / XLSX / PPTX.

Every format shares one visual system (brand palette + confidence colour coding)
and surfaces the full V3 evidence picture: an at-a-glance scorecard (evidence score,
confidence, consistency, verification), the PICO research question, the confidence-
scored comparison table, evidence-ranking + risk-of-bias, the comparative safety
matrix, conflict reconciliation, extracted trial data, and fully-linked references.
Nothing is invented here — exporters only reshape what the pipeline already verified.
"""

from __future__ import annotations

import io
import re
from datetime import UTC, datetime

from app.models.report import Export, Report

# (media_type, file extension) per export format.
EXPORT_MEDIA: dict[str, tuple[str, str]] = {
    "markdown": ("text/markdown; charset=utf-8", "md"),
    "pdf": ("application/pdf", "pdf"),
    "xlsx": (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "xlsx",
    ),
    "pptx": (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "pptx",
    ),
}

# ---- Brand system (shared across all formats) -------------------------------
BRAND = (13, 148, 136)       # teal-600 — primary
BRAND_DARK = (15, 23, 42)    # slate-900 — ink on light
ACCENT = (2, 132, 199)       # sky-600 — molecule B / accents
INK = (30, 41, 59)           # slate-800 — body text
MUTED = (100, 116, 139)      # slate-500 — secondary text
LIGHT = (241, 245, 249)      # slate-100 — panels
ZEBRA = (247, 249, 251)      # near-white — table striping
WHITE = (255, 255, 255)
_CONF = {
    "high": (22, 163, 74),
    "moderate": (202, 138, 4),
    "low": (234, 88, 12),
    "very_low": (220, 38, 38),
}


def _conf_rgb(level: str | None) -> tuple[int, int, int]:
    return _CONF.get((level or "").lower(), MUTED)


def _hex(rgb: tuple[int, int, int]) -> str:
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


_SLUG_RE = re.compile(r"[^a-z0-9]+")
# Map common non-latin-1 glyphs so the PDF core font (Helvetica) can render them.
_UNI = {
    "–": "-", "—": "-", "‘": "'", "’": "'", "“": '"',
    "”": '"', "…": "...", "≥": ">=", "≤": "<=", "×": "x",
    "−": "-", "±": "+/-", "α": "alpha", "β": "beta",
    "γ": "gamma", "µ": "u", "°": "deg", "→": "->", "‰": "0/00",
}


def _slug(report: Report) -> str:
    base = f"evidencecompare-{report.molecule_a}-vs-{report.molecule_b}"
    return _SLUG_RE.sub("-", base.lower()).strip("-") or "report"


def _ascii(s: str | None) -> str:
    if not s:
        return ""
    for k, v in _UNI.items():
        s = s.replace(k, v)
    return s.encode("latin-1", "replace").decode("latin-1")


def _claim_text(claim: object) -> str:
    return claim.get("text", "") if isinstance(claim, dict) else str(claim)


def _claim_cids(claim: object) -> list:
    return claim.get("citation_ids", []) if isinstance(claim, dict) else []


def _today() -> str:
    return datetime.now(UTC).strftime("%d %b %Y")


# ---- Shared data accessors (V3 JSON columns are nullable) -------------------
def _overall(report: Report) -> dict:
    return (report.evidence_scores or {}).get("overall", {}) if report.evidence_scores else {}


def _study_scores(report: Report) -> list[dict]:
    return (report.evidence_scores or {}).get("studies", []) if report.evidence_scores else []


def _pico(report: Report) -> dict | None:
    return (report.research_plan or {}).get("pico") if report.research_plan else None


def _safety_rows(report: Report) -> list[dict]:
    return (report.safety_matrix or {}).get("rows", []) if report.safety_matrix else []


def _scorecard(report: Report) -> list[tuple[str, str]]:
    """The at-a-glance metric tiles shared by PDF/XLSX/PPTX."""
    ov = _overall(report)
    verified = sum(1 for c in report.citations if c.verified)
    conf = str(ov.get("confidence", "n/a")).replace("_", " ")
    return [
        ("Evidence score", f"{ov.get('evidence_score', '-')}/100"),
        ("Confidence", conf),
        ("Consistency", f"{ov.get('consistency_score', '-')}/100"),
        ("Verified citations", str(verified or len(report.citations))),
        ("Studies scored", str(ov.get("n_studies", len(report.extractions)))),
    ]


def _cite_url(c) -> str:
    if c.doi:
        return f"https://doi.org/{c.doi}"
    if c.pmid:
        return f"https://pubmed.ncbi.nlm.nih.gov/{c.pmid}/"
    return ""


# =============================================================================
# Markdown
# =============================================================================
def _render_markdown(report: Report) -> str:
    a, b = report.molecule_a, report.molecule_b
    L: list[str] = []
    L.append(f"# EvidenceCompare AI — {a} vs {b}")
    L.append("")
    L.append(f"**Clinical topic:** {report.topic}  ")
    L.append(f"**Generated:** {_today()}  ")
    if report.model_synthesis:
        L.append(f"**Synthesis model:** {report.model_synthesis}  ")
    if report.freshness and report.freshness != "unknown":
        L.append(f"**Living evidence:** {report.freshness.replace('_', ' ')}  ")
    L.append("")

    # At a glance
    L.append("## At a glance")
    L.append("")
    L.append("| Metric | Value |")
    L.append("|---|---|")
    for label, value in _scorecard(report):
        L.append(f"| {label} | {value} |")
    L.append("")

    # PICO
    pico = _pico(report)
    if pico:
        L.append("## Research question (PICO)")
        L.append("")
        L.append(f"- **Population:** {pico.get('population', '-')}")
        L.append(f"- **Intervention:** {pico.get('intervention', a)}")
        L.append(f"- **Comparator:** {pico.get('comparator', b)}")
        L.append(f"- **Outcome:** {pico.get('outcome', report.topic)}")
        L.append("")

    if report.conflicts:
        L.append("## ⚠ Conflicting evidence")
        L.append("")
        for c in report.conflicts:
            L.append(f"- {c}")
        L.append("")

    if report.comparison_rows:
        L.append("## Side-by-side comparison")
        L.append("")
        L.append(f"| Attribute | {a} | {b} | Confidence | Rationale |")
        L.append("|---|---|---|---|---|")
        for row in report.comparison_rows:
            rat = (row.rationale or "").replace("|", "\\|")
            L.append(
                f"| {row.attribute} | {row.value_a} | {row.value_b} "
                f"| {row.confidence} | {rat} |"
            )
        L.append("")

    # Evidence ranking
    studies = _study_scores(report)
    if studies:
        L.append("## Evidence ranking")
        L.append("")
        L.append("| Ref | Study | Tier | Evidence | Risk of bias |")
        L.append("|---|---|---|---|---|")
        for s in studies:
            title = str(s.get("title", ""))[:60].replace("|", "\\|")
            L.append(
                f"| {s.get('ref_key')} | {title} | {s.get('tier_label')} "
                f"| {s.get('evidence_score')}/100 | {s.get('risk_of_bias')} |"
            )
        L.append("")

    # Safety matrix
    rows = _safety_rows(report)
    if rows:
        L.append("## Comparative safety matrix")
        L.append("")
        L.append(f"| Domain | {a} | {b} |")
        L.append("|---|---|---|")
        for r in rows:
            L.append(f"| {r['label']} | {r['a']['note']} | {r['b']['note']} |")
        L.append("")

    # Conflict reconciliation
    recon = report.reconciliation or {}
    if recon.get("explanations"):
        L.append("## Evidence reconciliation")
        L.append("")
        L.append(f"_{recon.get('summary', '')}_")
        L.append("")
        for ex in recon["explanations"]:
            L.append(f"- {ex.get('text', '')}")
        L.append("")

    ref_index = {c.ref_key: i + 1 for i, c in enumerate(report.citations)}

    for section in report.sections:
        L.append(f"## {section.title}")
        flag = " _(insufficient evidence)_" if section.insufficient_evidence else ""
        L.append(f"_Confidence: {section.confidence}{flag}_")
        L.append("")
        for claim in section.claims:
            cids = _claim_cids(claim)
            marks = "".join(f"[{ref_index.get(cid, '?')}]" for cid in cids)
            suffix = f" {marks}" if marks else " _(unsourced)_"
            L.append(f"- {_claim_text(claim)}{suffix}")
        L.append("")

    # Extracted trial data
    if report.extractions:
        L.append("## Extracted trial data")
        L.append("")
        L.append("| Ref | Study | Design | N | HR / RR | 95% CI | p |")
        L.append("|---|---|---|---|---|---|---|")
        for e in report.extractions:
            eff = e.hazard_ratio or e.relative_risk or "-"
            title = (e.title or "")[:50].replace("|", "\\|")
            L.append(
                f"| {e.ref_key} | {title} | {e.study_design or '-'} | {e.sample_size or '-'} "
                f"| {eff} | {e.confidence_interval or '-'} | {e.p_value or '-'} |"
            )
        L.append("")

    if report.citations:
        L.append("## References")
        L.append("")
        for i, c in enumerate(report.citations):
            bits = [c.title or "", c.source.upper()]
            if c.year:
                bits.append(str(c.year))
            if c.pmid:
                bits.append(f"PMID {c.pmid}")
            url = _cite_url(c)
            doi = f"[DOI {c.doi}]({url})" if c.doi and url else (c.doi or "")
            if doi:
                bits.append(doi)
            verified = " ✅" if c.verified else ""
            L.append(f"{i + 1}. {' · '.join(b for b in bits if b)}{verified}")
        L.append("")

    L.append("---")
    L.append(
        "> **Disclaimer.** Decision-support only. Not a diagnostic device or a "
        "substitute for professional clinical judgment."
    )
    return "\n".join(L)


# =============================================================================
# PDF (fpdf2) — branded cover, scorecards, styled tables
# =============================================================================
def _render_pdf(report: Report) -> bytes:
    from fpdf import FPDF  # type: ignore[import-untyped]
    from fpdf.enums import XPos, YPos  # type: ignore[import-untyped]
    from fpdf.fonts import FontFace  # type: ignore[import-untyped]

    a, b = _ascii(report.molecule_a), _ascii(report.molecule_b)

    class PDF(FPDF):
        def footer(self) -> None:
            self.set_y(-12)
            self.set_font("Helvetica", "I", 7)
            self.set_text_color(*MUTED)
            self.cell(
                0, 5,
                "EvidenceCompare AI  -  decision-support only, not a diagnostic device.",
                new_x=XPos.LMARGIN, new_y=YPos.TOP, align="L",
            )
            self.set_y(-12)
            self.cell(0, 5, f"Page {self.page_no()}", align="R")

    pdf = PDF(format="A4")
    pdf.set_auto_page_break(auto=True, margin=16)
    pdf.add_page()
    epw = pdf.epw  # effective page width

    def text(s: str, h: float = 5.0, size: float = 10, style: str = "",
             color: tuple[int, int, int] = INK) -> None:
        pdf.set_font("Helvetica", style, size)
        pdf.set_text_color(*color)
        pdf.multi_cell(0, h, _ascii(s), new_x=XPos.LMARGIN, new_y=YPos.NEXT)

    def h2(title: str) -> None:
        pdf.ln(3)
        y = pdf.get_y()
        pdf.set_fill_color(*BRAND)
        pdf.rect(pdf.l_margin, y, 2.5, 6, style="F")  # accent tab
        pdf.set_xy(pdf.l_margin + 4, y)
        pdf.set_font("Helvetica", "B", 13)
        pdf.set_text_color(*BRAND_DARK)
        pdf.cell(0, 6, _ascii(title), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.ln(1)
        # Reset the drawing fill to the zebra tone so table row-striping (which
        # picks up the current fill color) stays light, not the accent teal.
        pdf.set_fill_color(*ZEBRA)

    # ---- Cover band ----
    pdf.set_fill_color(*BRAND_DARK)
    pdf.rect(0, 0, pdf.w, 34, style="F")
    pdf.set_xy(pdf.l_margin, 8)
    pdf.set_font("Helvetica", "B", 20)
    pdf.set_text_color(*WHITE)
    pdf.cell(0, 9, f"{a}  vs  {b}", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    pdf.set_x(pdf.l_margin)
    pdf.set_font("Helvetica", "", 11)
    pdf.set_text_color(180, 190, 205)
    pdf.cell(0, 6, _ascii(f"{report.topic}   -   EvidenceCompare AI   -   {_today()}"))
    pdf.set_xy(pdf.l_margin, 40)
    pdf.set_text_color(*INK)

    # ---- Scorecard tiles ----
    tiles = _scorecard(report)
    gap = 4
    tw = (epw - gap * (len(tiles) - 1)) / len(tiles)
    x0, y0 = pdf.l_margin, pdf.get_y()
    for i, (label, value) in enumerate(tiles):
        x = x0 + i * (tw + gap)
        pdf.set_fill_color(*LIGHT)
        pdf.rect(x, y0, tw, 18, style="F")
        pdf.set_fill_color(*BRAND)
        pdf.rect(x, y0, tw, 1.6, style="F")  # top accent
        pdf.set_xy(x, y0 + 3)
        pdf.set_font("Helvetica", "B", 13)
        pdf.set_text_color(*BRAND_DARK)
        pdf.cell(tw, 7, _ascii(str(value)), align="C")
        pdf.set_xy(x, y0 + 11)
        pdf.set_font("Helvetica", "", 7.5)
        pdf.set_text_color(*MUTED)
        pdf.cell(tw, 4, _ascii(label.upper()), align="C")
    pdf.set_xy(pdf.l_margin, y0 + 22)

    # ---- PICO ----
    pico = _pico(report)
    if pico:
        h2("Research question (PICO)")
        for k, label in (("population", "Population"), ("intervention", "Intervention"),
                         ("comparator", "Comparator"), ("outcome", "Outcome")):
            pdf.set_font("Helvetica", "B", 9)
            pdf.set_text_color(*BRAND)
            pdf.cell(26, 5, f"{label}:", new_x=XPos.RIGHT, new_y=YPos.TOP)
            pdf.set_font("Helvetica", "", 9)
            pdf.set_text_color(*INK)
            pdf.multi_cell(0, 5, _ascii(str(pico.get(k, "-"))),
                           new_x=XPos.LMARGIN, new_y=YPos.NEXT)

    # ---- Conflicts banner ----
    if report.conflicts:
        h2("Conflicting evidence detected")
        for c in report.conflicts:
            text(f"- {c}", 5, 9, color=_CONF["low"])

    heading = FontFace(emphasis="BOLD", color=WHITE, fill_color=BRAND)

    # ---- Comparison table ----
    if report.comparison_rows:
        h2("Side-by-side comparison")
        pdf.set_font("Helvetica", "", 8.5)
        with pdf.table(
            col_widths=(30, 26, 26, 15),
            text_align=("LEFT", "LEFT", "LEFT", "CENTER"),
            headings_style=heading,
            cell_fill_color=ZEBRA,
            cell_fill_mode="ROWS",
            line_height=5,
            width=epw,
        ) as table:
            hdr = table.row()
            for c in ("Attribute", a, b, "Confidence"):
                hdr.cell(_ascii(c))
            for row in report.comparison_rows:
                r = table.row()
                r.cell(_ascii(row.attribute))
                r.cell(_ascii(row.value_a))
                r.cell(_ascii(row.value_b))
                r.cell(_ascii(row.confidence),
                       style=FontFace(color=WHITE, fill_color=_conf_rgb(row.confidence)))

    # ---- Evidence ranking ----
    studies = _study_scores(report)
    if studies:
        h2("Evidence ranking")
        ov = _overall(report)
        text(
            f"Aggregate evidence {ov.get('evidence_score', '-')}/100, "
            f"{str(ov.get('confidence', '')).replace('_', ' ')} certainty; "
            f"consistency {ov.get('consistency_score', '-')}/100; "
            f"overall risk of bias {ov.get('risk_of_bias', '-')}.",
            5, 9, color=MUTED,
        )
        pdf.set_font("Helvetica", "", 8.5)
        with pdf.table(
            col_widths=(10, 44, 20, 14, 15),
            text_align=("CENTER", "LEFT", "LEFT", "CENTER", "CENTER"),
            headings_style=heading, cell_fill_color=ZEBRA, cell_fill_mode="ROWS",
            line_height=5, width=epw,
        ) as table:
            hdr = table.row()
            for c in ("Ref", "Study", "Tier", "Evid.", "Risk"):
                hdr.cell(_ascii(c))
            for s in studies[:12]:
                r = table.row()
                r.cell(_ascii(str(s.get("ref_key", ""))))
                r.cell(_ascii(str(s.get("title", ""))[:70]))
                r.cell(_ascii(str(s.get("tier_label", ""))))
                r.cell(_ascii(f"{s.get('evidence_score', '')}"))
                r.cell(_ascii(str(s.get("risk_of_bias", ""))))

    # ---- Safety matrix ----
    rows = _safety_rows(report)
    if rows:
        h2("Comparative safety matrix")
        pdf.set_font("Helvetica", "", 8.5)
        with pdf.table(
            col_widths=(28, 34, 34),
            text_align=("LEFT", "LEFT", "LEFT"),
            headings_style=heading, cell_fill_color=ZEBRA, cell_fill_mode="ROWS",
            line_height=5, width=epw,
        ) as table:
            hdr = table.row()
            for c in ("Domain", a, b):
                hdr.cell(_ascii(c))
            for srow in rows:
                tr = table.row()
                tr.cell(_ascii(srow["label"]))
                tr.cell(_ascii(srow["a"]["note"]))
                tr.cell(_ascii(srow["b"]["note"]))

    # ---- Reconciliation ----
    recon = report.reconciliation or {}
    if recon.get("explanations"):
        h2("Evidence reconciliation")
        text(recon.get("summary", ""), 5, 9, "I", MUTED)
        for ex in recon["explanations"]:
            text(f"- {ex.get('text', '')}", 5, 9)

    # ---- Narrative sections ----
    ref_index = {c.ref_key: i + 1 for i, c in enumerate(report.citations)}
    for section in report.sections:
        conf = _ascii(section.confidence)
        h2(section.title)
        pdf.set_font("Helvetica", "B", 8)
        pdf.set_text_color(*WHITE)
        pdf.set_fill_color(*_conf_rgb(section.confidence))
        pdf.cell(pdf.get_string_width(conf) + 6, 5, conf, align="C", fill=True,
                 new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.ln(1)
        for claim in section.claims:
            marks = "".join(f"[{ref_index.get(c, '?')}]" for c in _claim_cids(claim))
            text(f"- {_claim_text(claim)} {marks}".rstrip(), 5, 9.5)

    # ---- Extracted trial data ----
    if report.extractions:
        h2("Extracted trial data")
        pdf.set_font("Helvetica", "", 8)
        with pdf.table(
            col_widths=(10, 40, 16, 10, 16, 16),
            text_align=("CENTER", "LEFT", "LEFT", "CENTER", "LEFT", "LEFT"),
            headings_style=heading, cell_fill_color=ZEBRA, cell_fill_mode="ROWS",
            line_height=4.6, width=epw,
        ) as table:
            hdr = table.row()
            for c in ("Ref", "Study", "Design", "N", "HR / RR", "95% CI"):
                hdr.cell(_ascii(c))
            for e in report.extractions:
                eff = e.hazard_ratio or e.relative_risk or "-"
                r = table.row()
                r.cell(_ascii(e.ref_key))
                r.cell(_ascii((e.title or "")[:60]))
                r.cell(_ascii(e.study_design or "-"))
                r.cell(_ascii(str(e.sample_size or "-")))
                r.cell(_ascii(eff))
                r.cell(_ascii(e.confidence_interval or "-"))

    # ---- References ----
    if report.citations:
        h2("References")
        pdf.set_font("Helvetica", "", 8.5)
        for i, c in enumerate(report.citations):
            bits = [c.title or "", c.source.upper()]
            if c.year:
                bits.append(str(c.year))
            if c.pmid:
                bits.append(f"PMID {c.pmid}")
            if c.doi:
                bits.append(f"DOI {c.doi}")
            tag = "  [verified]" if c.verified else ""
            text(f"{i + 1}. " + "  |  ".join(bits) + tag, 4.6, 8.5, color=INK)

    return bytes(pdf.output())


# =============================================================================
# XLSX (openpyxl) — styled multi-sheet workbook
# =============================================================================
def _render_xlsx(report: Report) -> bytes:
    from openpyxl import Workbook  # type: ignore[import-untyped]
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
    from openpyxl.utils import get_column_letter

    a, b = report.molecule_a, report.molecule_b
    thin = Side(style="thin", color="D0D7DE")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    head_fill = PatternFill("solid", fgColor=_hex(BRAND))
    head_font = Font(bold=True, color="FFFFFF")
    title_font = Font(bold=True, size=15, color=_hex(BRAND_DARK))
    wrap = Alignment(wrap_text=True, vertical="top")

    wb = Workbook()

    def sheet(name: str, header: list[str], data: list[list], title: str | None = None):
        ws = wb.create_sheet(name)
        if title:
            ws.append([title])
            ws["A1"].font = title_font
            ws.append([])
        hrow = ws.max_row + 1
        ws.append(header)
        for cell in ws[hrow]:
            cell.fill, cell.font, cell.border = head_fill, head_font, border
            cell.alignment = Alignment(vertical="center")
        for row in data:
            ws.append(row)
        # borders + wrap on body
        for r in ws.iter_rows(min_row=hrow + 1, max_row=ws.max_row, max_col=len(header)):
            for cell in r:
                cell.border, cell.alignment = border, wrap
        # auto column widths
        for col in range(1, len(header) + 1):
            letter = get_column_letter(col)
            longest = max(
                (len(str(ws.cell(row=r, column=col).value or ""))
                 for r in range(hrow, ws.max_row + 1)),
                default=10,
            )
            ws.column_dimensions[letter].width = min(max(longest + 2, 12), 60)
        ws.freeze_panes = ws.cell(row=hrow + 1, column=1)
        return ws

    # Overview
    ov = _overall(report)
    pico = _pico(report) or {}
    ws = wb.active
    ws.title = "Overview"
    ws["A1"] = f"EvidenceCompare AI — {a} vs {b}"
    ws["A1"].font = title_font
    meta = [
        ("Clinical topic", report.topic),
        ("Generated", _today()),
        ("Synthesis model", report.model_synthesis or ""),
        ("Living evidence", (report.freshness or "").replace("_", " ")),
        ("", ""),
        ("Evidence score", f"{ov.get('evidence_score', '-')}/100"),
        ("Confidence", str(ov.get("confidence", "-")).replace("_", " ")),
        ("Consistency", f"{ov.get('consistency_score', '-')}/100"),
        ("Overall risk of bias", ov.get("risk_of_bias", "-")),
        ("Verified citations", sum(1 for c in report.citations if c.verified)),
        ("", ""),
        ("PICO — Population", pico.get("population", "")),
        ("PICO — Intervention", pico.get("intervention", a)),
        ("PICO — Comparator", pico.get("comparator", b)),
        ("PICO — Outcome", pico.get("outcome", report.topic)),
    ]
    for r, (k, v) in enumerate(meta, start=3):
        ws.cell(row=r, column=1, value=k).font = Font(bold=True, color=_hex(MUTED))
        ws.cell(row=r, column=2, value=v)
    ws.column_dimensions["A"].width = 24
    ws.column_dimensions["B"].width = 70

    # Comparison (with confidence colour coding)
    cw = sheet(
        "Comparison", ["Attribute", a, b, "Confidence", "Rationale"],
        [[r.attribute, r.value_a, r.value_b, r.confidence, r.rationale or ""]
         for r in report.comparison_rows],
    )
    for row in cw.iter_rows(min_row=2, max_col=4):
        conf_cell = row[3]
        if conf_cell.value:
            conf_cell.fill = PatternFill("solid", fgColor=_hex(_conf_rgb(str(conf_cell.value))))
            conf_cell.font = Font(bold=True, color="FFFFFF")

    # Evidence scores
    studies = _study_scores(report)
    if studies:
        sheet(
            "Evidence Scores",
            ["Ref", "Study", "Tier", "Evidence /100", "Quality /100",
             "Risk of bias", "Publication /100"],
            [[s.get("ref_key"), s.get("title"), s.get("tier_label"),
              s.get("evidence_score"), s.get("quality_score"),
              s.get("risk_of_bias"), s.get("publication_quality")] for s in studies],
        )

    # Safety matrix
    rows = _safety_rows(report)
    if rows:
        sheet(
            "Safety Matrix", ["Domain", f"{a}", f"{b}"],
            [[r["label"], r["a"]["note"], r["b"]["note"]] for r in rows],
        )

    # Extractions
    if report.extractions:
        sheet(
            "Extractions",
            ["Ref", "Title", "Design", "Population", "Intervention", "Comparator",
             "N", "HR", "RR", "95% CI", "p", "Adverse events"],
            [[e.ref_key, e.title, e.study_design or "", e.population or "",
              e.intervention or "", e.comparator or "", e.sample_size or "",
              e.hazard_ratio or "", e.relative_risk or "", e.confidence_interval or "",
              e.p_value or "", ", ".join(e.adverse_events or [])] for e in report.extractions],
        )

    # References
    sheet(
        "References", ["#", "Title", "Source", "Year", "PMID", "DOI", "URL", "Verified"],
        [[i + 1, c.title, c.source, c.year or "", c.pmid or "", c.doi or "",
          _cite_url(c), "yes" if c.verified else "no"]
         for i, c in enumerate(report.citations)],
    )

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# =============================================================================
# PPTX (python-pptx) — branded 16:9 deck
# =============================================================================
def _render_pptx(report: Report) -> bytes:
    from pptx import Presentation  # type: ignore[import-untyped]
    from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
    from pptx.enum.text import PP_ALIGN  # type: ignore[import-untyped]
    from pptx.util import Inches, Pt  # type: ignore[import-untyped]

    a, b = report.molecule_a, report.molecule_b
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    BLANK = prs.slide_layouts[6]

    def rgb(t: tuple[int, int, int]) -> RGBColor:
        return RGBColor(*t)

    def rect(slide, x, y, w, h, color):
        from pptx.enum.shapes import MSO_SHAPE  # type: ignore[import-untyped]
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(color)
        shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def textbox(slide, x, y, w, h, text, size, *, color=INK, bold=False,
                align=PP_ALIGN.LEFT, font="Calibri"):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        run.font.color.rgb = rgb(color)
        return tb

    def content_slide(title: str):
        slide = prs.slides.add_slide(BLANK)
        rect(slide, 0, 0, SW, Inches(0.16), BRAND)  # top accent bar
        textbox(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.8),
                title, 26, color=BRAND_DARK, bold=True)
        return slide

    def bullets(slide, x, y, w, h, items, size=15):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"•  {it}"
            p.space_after = Pt(6)
            for run in p.runs:
                run.font.size = Pt(size)
                run.font.color.rgb = rgb(INK)
                run.font.name = "Calibri"

    def table_slide(title, headers, data, widths):
        slide = content_slide(title)
        rows, cols = len(data) + 1, len(headers)
        left, top = Inches(0.6), Inches(1.4)
        width = SW - Inches(1.2)
        height = min(SH - Inches(1.8), Inches(0.4) * rows)
        gt = slide.shapes.add_table(rows, cols, left, top, width, height).table
        for j, w in enumerate(widths):
            gt.columns[j].width = Inches(w)
        for j, htext in enumerate(headers):
            cell = gt.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(BRAND)
            p = cell.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = str(htext)
            r.font.bold = True
            r.font.color.rgb = rgb(WHITE)
            r.font.size = Pt(11)
        for i, rowdata in enumerate(data, start=1):
            for j, val in enumerate(rowdata):
                cell = gt.cell(i, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(WHITE if i % 2 else ZEBRA)
                p = cell.text_frame.paragraphs[0]
                r = p.add_run()
                r.text = str(val)
                r.font.size = Pt(10)
                r.font.color.rgb = rgb(INK)
        return slide

    # ---- Title slide ----
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, BRAND_DARK)
    rect(s, 0, Inches(4.55), SW, Inches(0.06), BRAND)
    textbox(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.4),
            f"{a}  vs  {b}", 46, color=WHITE, bold=True)
    textbox(s, Inches(0.95), Inches(3.9), Inches(11.5), Inches(0.7),
            report.topic, 22, color=(180, 190, 205))
    textbox(s, Inches(0.95), Inches(4.75), Inches(11.5), Inches(0.5),
            f"EvidenceCompare AI   ·   {_today()}", 13, color=(148, 163, 184))

    # ---- At a glance scorecards ----
    s = content_slide("At a glance")
    tiles = _scorecard(report)
    tw = (SW - Inches(1.2) - Inches(0.3) * (len(tiles) - 1)) / len(tiles)
    for i, (label, value) in enumerate(tiles):
        x = Inches(0.6) + i * (tw + Inches(0.3))
        card = rect(s, x, Inches(1.7), tw, Inches(1.7), LIGHT)
        card.shadow.inherit = False
        rect(s, x, Inches(1.7), tw, Inches(0.12), BRAND)
        textbox(s, x, Inches(2.1), tw, Inches(0.8), str(value), 26,
                color=BRAND_DARK, bold=True, align=PP_ALIGN.CENTER)
        textbox(s, x, Inches(2.95), tw, Inches(0.4), label.upper(), 10,
                color=MUTED, align=PP_ALIGN.CENTER)
    pico = _pico(report)
    if pico:
        bullets(s, Inches(0.6), Inches(3.9), SW - Inches(1.2), Inches(3),
                [f"Population: {pico.get('population', '-')}",
                 f"Intervention: {pico.get('intervention', a)}",
                 f"Comparator: {pico.get('comparator', b)}",
                 f"Outcome: {pico.get('outcome', report.topic)}"], size=14)

    # ---- Key clinical findings ----
    summary = [
        _claim_text(cl) for sec in report.sections
        if sec.layer == "clinical_summary" for cl in sec.claims
    ]
    if summary:
        s = content_slide("Key clinical findings")
        bullets(s, Inches(0.6), Inches(1.5), SW - Inches(1.2), Inches(5.5),
                summary[:8])

    # ---- Comparison table ----
    if report.comparison_rows:
        table_slide(
            "Side-by-side comparison",
            ["Attribute", a, b, "Conf."],
            [[r.attribute, r.value_a[:60], r.value_b[:60], r.confidence]
             for r in report.comparison_rows[:8]],
            [3.6, 3.6, 3.6, 1.3],
        )

    # ---- Evidence ranking ----
    studies = _study_scores(report)
    if studies:
        table_slide(
            "Evidence ranking",
            ["Ref", "Study", "Tier", "Evid.", "Risk"],
            [[s2.get("ref_key"), str(s2.get("title", ""))[:48], s2.get("tier_label"),
              f"{s2.get('evidence_score')}/100", s2.get("risk_of_bias")]
             for s2 in studies[:9]],
            [1.0, 5.3, 2.2, 1.4, 1.7],
        )

    # ---- Safety matrix ----
    rows = _safety_rows(report)
    if rows:
        table_slide(
            "Comparative safety matrix",
            ["Domain", a, b],
            [[r["label"], r["a"]["note"][:55], r["b"]["note"][:55]] for r in rows],
            [3.1, 4.5, 4.5],
        )

    # ---- Reconciliation ----
    recon = report.reconciliation or {}
    if recon.get("explanations"):
        s = content_slide("Evidence reconciliation")
        bullets(s, Inches(0.6), Inches(1.5), SW - Inches(1.2), Inches(5),
                [recon.get("summary", "")] + [e.get("text", "") for e in recon["explanations"]])

    # ---- References ----
    if report.citations:
        s = content_slide("References")
        bullets(s, Inches(0.6), Inches(1.4), SW - Inches(1.2), Inches(5.7),
                [f"{i + 1}. {c.title} ({c.source.upper()}"
                 f"{', ' + str(c.year) if c.year else ''})"
                 for i, c in enumerate(report.citations)][:14], size=12)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def render_export(report: Report, fmt: str) -> tuple[bytes, str, str]:
    """Render a report to downloadable bytes. Returns (data, media_type, filename)."""
    media, ext = EXPORT_MEDIA[fmt]
    if fmt == "markdown":
        data = _render_markdown(report).encode("utf-8")
    elif fmt == "pdf":
        data = _render_pdf(report)
    elif fmt == "xlsx":
        data = _render_xlsx(report)
    elif fmt == "pptx":
        data = _render_pptx(report)
    else:  # pragma: no cover - schema restricts fmt
        raise ValueError(f"unsupported format: {fmt}")
    return data, media, f"{_slug(report)}.{ext}"


def build_export(report: Report, fmt: str) -> Export:
    """Record an export request. Markdown content is inlined; binary formats are
    generated on demand by the download endpoint via `render_export`."""
    content = _render_markdown(report) if fmt == "markdown" else None
    _, ext = EXPORT_MEDIA[fmt]
    return Export(
        report_id=report.id,
        format=fmt,
        status="ready",
        object_url=f"/api/v1/reports/{report.id}/download?format={fmt}",
        content=content,
    )
